import streamlit as st
import os
import random
import json
import zipfile
import io
import docx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement  # Stable XML module [1.1.4]
from pptx.oxml.ns import qn
import google.generativeai as genai

# ==========================================
# 1. GLOBAL PREMIUM THEME COLORS (Wow Factor)
# ==========================================
COLOR_ROYAL_BLUE = RGBColor(10, 30, 80)     # Deep Royal Blue (Main Academic)
COLOR_TEAL = RGBColor(0, 130, 140)          # Tech/Science Teal
COLOR_EMERALD = RGBColor(40, 160, 80)       # Success/Vibrant Green
COLOR_RUST = RGBColor(220, 100, 40)         # Alert/Warning Rust Orange
COLOR_GOLD = RGBColor(240, 180, 40)         # Highlight Gold/Amber
COLOR_LIGHT_GRAY = RGBColor(246, 248, 250)   # Clean Slide Background Texture [1]
COLOR_WHITE = RGBColor(255, 255, 255)       # Solid White Cards
COLOR_TEXT_DARK = RGBColor(40, 40, 40)      # Readable body text
COLOR_ACCENT_GREEN = RGBColor(40, 167, 69)  # High-impact green stats

# DYNAMIC CONTEXTUAL BACKGROUND COLORS [1]
BG_TITLE = RGBColor(15, 34, 64)               # Dark Blue background for Title Slide [1]
BG_PROBLEM = RGBColor(253, 240, 235)          # Soft Alert Orange for Problem Statement [1]
BG_OBJECTIVE = RGBColor(235, 247, 245)        # Soft Mint Teal for Objectives [1]
BG_METHODOLOGY = RGBColor(255, 255, 255)      # Crisp White (to contrast colorful cards) [1]
BG_RESULTS = RGBColor(240, 244, 252)          # Pale Ice Blue for Results/Data [1]
BG_CONCLUSION = RGBColor(235, 248, 238)       # Soft Emerald Green for Conclusions [1]
BG_DEFAULT = RGBColor(246, 248, 250)          # Classic Light Gray

# ==========================================
# 2. PAGE SETUP & SECURITY
# ==========================================
st.set_page_config(page_title="M.Sc. Premium Thesis PPT Generator", page_icon="🎓", layout="centered")

api_key = st.secrets.get("GEMINI_API_KEY") or os.environ.get("GEMINI_API_KEY")

if api_key:
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-2.5-flash')
else:
    st.warning("⚠️ API Key not detected. Please set GEMINI_API_KEY in Streamlit Secrets.")

# ==========================================
# 3. FILE READERS
# ==========================================
def read_word_file(file_obj):
    doc = docx.Document(file_obj)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def read_ppt_file(file_obj):
    prs = Presentation(file_obj)
    slide_texts = []
    for slide in prs.slides:
        text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        if text.strip(): slide_texts.append(text)
    return slide_texts

# ==========================================
# 4. EXTRA-VISUAL-FILTER ACADEMIC AI ENGINE
# ==========================================
def ai_process_word_advanced(full_text):
    prompt = f"""
    You are an expert Academic Presentation Assistant. Your goal is to structure this text into HIGHLY VISUAL, PREMIUM presentation slides with a 'WOW' factor.
    
    STRICT TEXT LIMITS:
    - Never write long paragraphs or sentences on slides.
    - Each slide bullet point must be strictly restricted to 3-5 words maximum.
    - Focus strictly on dynamic visual representation rather than plain text.
    - Put all extensive explanations in the 'script' (Speaker Notes) so the presenter can speak for 1 minute.

    Available Layout Types:
    1. 'process': Use for workflows, methodologies, or step-by-step procedures. (Horizontal process cards with badges)
    2. 'table': Use for comparative parameter matrix, structured comparison data.
    3. 'comparison': Use for contrasting two concepts, Pros vs Cons, Old vs New. (Side-by-side split cards)
    4. 'stat': Use if there's a strong numeric fact, percentage, metric, or major statistic (e.g., '95% Efficiency', '500+ Participants').
    5. 'bullets': Use only if none of the above visual layouts match.

    Return ONLY a valid JSON object in this format (strictly no markdown formatting outside the JSON, no extra text):
    {{
        "category": "thesis",
        "recommended_source": "slidesgo | presentationgo | slidescarnival",
        "mermaid_diagram": "graph TD\\n    A[Step 1] --> B[Step 2]",
        "slides": [
            {{
                "title": "Slide Title",
                "layout_type": "process | table | comparison | stat | bullets",
                "content": {{
                    "bullets": ["Keyword-focused bullet 1", "Keyword-focused bullet 2"],
                    "table": {{
                        "headers": ["Parameter", "Detail A", "Detail B"],
                        "rows": [["Row 1 Col 1", "Row 1 Col 2", "Row 1 Col 3"]]
                    }},
                    "process": [
                        {{"step": "1", "title": "Phase 1", "desc": "Phase 1 summary"}}
                    ],
                    "comparison": {{
                        "left_title": "Side A Title",
                        "left_content": ["Short point A1", "Short point A2"],
                        "right_title": "Side B Title",
                        "right_content": ["Short point B1", "Short point B2"]
                    }},
                    "stat": {{
                        "number": "95%",
                        "label": "Metric Description (max 5 words)"
                    }}
                }},
                "script": "A detailed, professional 1-minute presenter speech explaining the slide content..."
            }}
        ]
    }}

    Academic Text:
    {full_text}
    """
    response = model.generate_content(prompt)
    clean_json = response.text.strip().replace('```json', '').replace('```', '')
    return json.loads(clean_json)

def ai_process_ppt_advanced(slide_texts):
    packaged_text = ""
    for idx, text in enumerate(slide_texts):
        packaged_text += f"\n--- ROUGH SLIDE {idx+1} ---\n{text}\n"
        
    prompt = f"""
    You are an expert Academic Presentation Assistant. Process this rough slide deck. Maintain 1:1 slide order ({len(slide_texts)} slides).
    Strictly summarize the text into a premium M.Sc. Defense visual design.
    
    STRICT 6x6 RULE:
    - Maximum 6 bullet lines, maximum 6 words per line [4]. Prioritize tables and flowcharts.
    - Put all speech explanations in the 'script' (Speaker Notes) so the presenter can speak for 1 minute [1].
    - Generate a beautiful MERMAID.JS flowchart code representing the methodology workflow.

    For each slide, choose the most appropriate layout_type: 'title', 'background', 'problem_statement', 'objective', 'methodology', 'results', or 'conclusion'.
    
    Return ONLY a valid JSON object in this format (strictly no markdown formatting outside the JSON, no extra text):
    {{
        "category": "thesis",
        "recommended_source": "slidesgo | presentationgo | slidescarnival",
        "mermaid_diagram": "graph TD\\n    A[Step 1] --> B[Step 2]",
        "slides": [
            {{
                "title": "Slide Title",
                "layout_type": "process | table | comparison | stat | bullets",
                "content": {{
                    "bullets": ["Keyword point 1", "Keyword point 2"],
                    "table": {{
                        "headers": ["Parameter", "Detail A", "Detail B"],
                        "rows": [["Row 1 Col 1", "Row 1 Col 2", "Row 1 Col 3"]]
                    }},
                    "process": [
                        {{"step": "1", "title": "Phase 1", "desc": "Phase 1 summary"}}
                    ],
                    "comparison": {{
                        "left_title": "Side A",
                        "left_content": ["Short point A1"],
                        "right_title": "Side B",
                        "right_content": ["Short point B1"]
                    }},
                    "stat": {{
                        "number": "85%",
                        "label": "Metric Description"
                    }}
                }},
                "script": "A detailed 1-minute presenter speech explaining the slide..."
            }}
        ]
    }}

    Rough Slides Data:
    {packaged_text}
    """
    response = model.generate_content(prompt)
    clean_json = response.text.strip().replace('```json', '').replace('```', '')
    return json.loads(clean_json)

# ==========================================
# 5. JSON ROBUST NORMALIZER
# ==========================================
def normalize_ai_data(ai_data):
    normalized = {"category": "lecture", "recommended_source": "slidesgo", "mermaid_diagram": "graph TD\n    A[Start] --> B[Process]", "slides": []}
    if not ai_data: return normalized
        
    if isinstance(ai_data, list):
        raw_slides = ai_data
    elif isinstance(ai_data, dict):
        normalized["category"] = ai_data.get("category", "lecture")
        normalized["recommended_source"] = ai_data.get("recommended_source", "slidesgo")
        normalized["mermaid_diagram"] = ai_data.get("mermaid_diagram", "graph TD\n    A[Start] --> B[Process]")
        raw_slides = ai_data.get("slides", [])
        if isinstance(raw_slides, dict):
            raw_slides = [raw_slides]
    else:
        raw_slides = []
        
    for item in raw_slides:
        if isinstance(item, list):
            if len(item) > 0: item = item[0]
            else: continue
                
        if isinstance(item, dict):
            slide = {
                "title": item.get("title", "Academic Slide"),
                "layout_type": item.get("layout_type", "bullets"),
                "content": item.get("content", {}),
                "script": item.get("script", "No speech generated.")
            }
            normalized["slides"].append(slide)
            
    return normalized

# ==========================================
# 6. PROGRAMMATIC TRANSITION & TIMING TRIGGER (XML)
# ==========================================
def set_slide_transition_and_timing(slide, transition_type="fade", auto_advance_sec=60):
    sLD = slide._element
    transition = sLD.find(qn('p:transition'))
    if transition is None:
        transition = OxmlElement('p:transition')
        sLD.append(transition)
        
    for child in list(transition):
        transition.remove(child)
        
    if transition_type == "fade":
        effect = OxmlElement('p:fade')
        transition.append(effect)
    elif transition_type == "push":
        effect = OxmlElement('p:push')
        effect.set('dir', 'lt')
        transition.append(effect)
    elif transition_type == "wipe":
        effect = OxmlElement('p:wipe')
        effect.set('dir', 'lt')
        transition.append(effect)
        
    if auto_advance_sec > 0:
        transition.set('advClick', '0')
        transition.set('advTm', str(auto_advance_sec * 1000))

# ==========================================
# 7. PREMIUM WOW-FACTOR VISUAL DRAWING ENGINE
# ==========================================
def apply_slide_branding_texture(slide, layout_type):
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_ROYAL_BLUE if layout_type != "title" else COLOR_TEAL
    top_bar.line.fill.background()

def draw_visual_thesis_slide(slide, slide_data):
    layout_type = slide_data.get("layout_type", "bullets")
    content = slide_data.get("content", {})
    
    # Contextual background selection [1]
    bg_color = BG_DEFAULT
    if layout_type == "title": bg_color = BG_TITLE [1]
    elif layout_type == "problem_statement": bg_color = BG_PROBLEM [1]
    elif layout_type == "objective": bg_color = BG_OBJECTIVE [1]
    elif layout_type == "methodology": bg_color = BG_METHODOLOGY [1]
    elif layout_type == "results": bg_color = BG_RESULTS [1]
    elif layout_type == "conclusion": bg_color = BG_CONCLUSION [1]

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    # Apply top branding bar
    apply_slide_branding_texture(slide, layout_type)
    
    # Title Setup
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get('title', "Academic Slide").upper()
        p = slide.shapes.title.text_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = COLOR_ROYAL_BLUE if layout_type != "title" else COLOR_WHITE

    # Clean default placeholders
    for shape in list(slide.shapes):
        if shape.is_placeholder and shape.placeholder_format.idx == 1:
            if layout_type in ["title", "problem_statement", "objective", "methodology", "results", "conclusion"]:
                sp = shape._element
                sp.getparent().remove(sp)

    # LAYOUT 1: TITLE SLIDE
    if layout_type == "title":
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(-0.5), Inches(0), Inches(4.5), Inches(7.5))
        panel.fill.solid()
        panel.fill.fore_color.rgb = COLOR_TEAL
        panel.line.fill.background()
        
        t_box = slide.shapes.add_textbox(Inches(4.8), Inches(2.2), Inches(8.0), Inches(3.5))
        tf = t_box.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = slide_data.get('title', "M.Sc. Thesis Defense").upper()
        p_t.font.bold = True
        p_t.font.size = Pt(40)
        p_t.font.color.rgb = COLOR_WHITE
        
        p_sub = tf.add_paragraph()
        p_sub.text = "\nM.Sc. Thesis Defense Presentation\nSchool of Science"
        p_sub.font.size = Pt(18)
        p_sub.font.color.rgb = COLOR_TEAL

    # LAYOUT 2: PROBLEM STATEMENT (Rust Orange card)
    elif layout_type == "problem_statement":
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.33), Inches(4.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_RUST
        card.line.width = Pt(2.5)
        
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.33), Inches(0.25))
        band.fill.solid()
        band.fill.fore_color.rgb = COLOR_RUST
        band.line.fill.background()
        
        tf = card.text_frame
        tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = "\n⚠️ KEY RESEARCH PROBLEM"
        p_title.font.bold = True
        p_title.font.size = Pt(18)
        p_title.font.color.rgb = COLOR_RUST
        
        bullets = content.get("bullets", ["Critical research problem statement."])
        for bullet in bullets[:4]:
            p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT_DARK
            p.space_after = Pt(10)

    # LAYOUT 3: OBJECTIVES (Target checklist cards)
    elif layout_type == "objective":
        bullets = content.get("bullets", ["Research Objective."])
        start_y = Inches(2.2)
        box_height = Inches(0.8)
        
        for idx, bullet in enumerate(bullets[:4]):
            y = start_y + idx * Inches(1.0)
            
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.6), Inches(0.6))
            badge.fill.solid()
            badge.fill.fore_color.rgb = COLOR_EMERALD
            badge.line.fill.background()
            badge.text_frame.text = f"O{idx+1}"
            p_b = badge.text_frame.paragraphs[0]
            p_b.font.bold = True
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = COLOR_WHITE
            p_b.alignment = PP_ALIGN.CENTER
            
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), y, Inches(9.5), box_height)
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_WHITE
            card.line.color.rgb = COLOR_TEAL
            
            tf = card.text_frame
            tf.word_wrap = True
            p_c = tf.paragraphs[0]
            p_c.text = bullet
            p_c.font.size = Pt(14)
            p_c.font.bold = True
            p_c.font.color.rgb = COLOR_TEXT_DARK

    # LAYOUT 4: METHODOLOGY (Flowchart process cards - Multi-Color Theme)
    elif layout_type == "methodology":
        steps = content.get("process", [])[:4]
        if not steps and "bullets" in content:
            steps = [{"step": str(i+1), "title": f"Step {i+1}", "desc": bullet} for i, bullet in enumerate(content["bullets"][:4])]
            
        num_steps = len(steps)
        if num_steps > 0:
            start_x = Inches(1.0)
            y = Inches(2.5)
            box_width = Inches((11.33 - (num_steps - 1) * 0.4) / num_steps)
            box_height = Inches(3.2)
            vibrant_colors = [COLOR_TEAL, COLOR_ROYAL_BLUE, COLOR_GOLD, COLOR_EMERALD]
            
            for idx, step_data in enumerate(steps):
                x = start_x + idx * (box_width + Inches(0.4))
                step_color = vibrant_colors[idx % len(vibrant_colors)]
                
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
                card.fill.solid()
                card.fill.fore_color.rgb = COLOR_WHITE
                card.line.color.rgb = step_color
                card.line.width = Pt(1.5)
                
                header_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, Inches(0.45))
                header_bar.fill.solid()
                header_bar.fill.fore_color.rgb = step_color
                header_bar.line.fill.background()
                
                badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + box_width/2 - Inches(0.3), y - Inches(0.3), Inches(0.6), Inches(0.6))
                badge.fill.solid()
                badge.fill.fore_color.rgb = step_color
                badge.line.fill.background()
                badge.text_frame.text = str(idx + 1)
                p_b = badge.text_frame.paragraphs[0]
                p_b.font.bold = True
                p_b.font.size = Pt(14)
                p_b.font.color.rgb = COLOR_WHITE
                p_b.alignment = PP_ALIGN.CENTER
                
                tf = card.text_frame
                tf.word_wrap = True
                p_space = tf.paragraphs[0]
                p_space.text = "\n"
                
                p_title = tf.add_paragraph()
                p_title.text = step_data.get('title', "").upper()
                p_title.font.bold = True
                p_title.font.size = Pt(14)
                p_title.font.color.rgb = step_color
                p_title.alignment = PP_ALIGN.CENTER
                
                p_desc = tf.add_paragraph()
                p_desc.text = f"\n{step_data.get('desc', '')}"
                p_desc.font.size = Pt(11)
                p_desc.font.color.rgb = COLOR_TEXT_DARK
                p_desc.alignment = PP_ALIGN.CENTER

                if idx < num_steps - 1:
                    arrow_x = x + box_width + Inches(0.05)
                    arrow_y = y + Inches(1.4)
                    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.3), Inches(0.4))
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = step_color
                    arrow.line.fill.background()

    # LAYOUT 5: RESULTS (Sleek Data Matrix table / High-Impact metrics)
    elif layout_type == "results":
        if "table" in content:
            table_data = content["table"]
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            num_cols = len(headers)
            num_rows = len(rows) + 1
            
            if num_cols > 0 and num_rows > 1:
                left = Inches(1.5)
                top = Inches(2.2)
                width = Inches(10.33)
                height = Inches(0.6 * num_rows)
                
                table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
                table = table_shape.table
                
                for col_idx, header_text in enumerate(headers):
                    cell = table.cell(0, col_idx)
                    cell.text = str(header_text).upper()
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_ROYAL_BLUE
                    p = cell.text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.color.rgb = COLOR_WHITE
                    p.font.size = Pt(13)
                    p.alignment = PP_ALIGN.CENTER
                    
                for row_idx, row_values in enumerate(rows):
                    for col_idx, val in enumerate(row_values):
                        if col_idx < num_cols:
                            cell = table.cell(row_idx + 1, col_idx)
                            cell.text = str(val)
                            cell.fill.solid()
                            if row_idx % 2 == 0:
                                cell.fill.fore_color.rgb = COLOR_WHITE
                            else:
                                cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY
                            p = cell.text_frame.paragraphs[0]
                            p.font.size = Pt(11)
                            p.font.color.rgb = COLOR_TEXT_DARK
                            p.alignment = PP_ALIGN.LEFT
        else:
            stat_data = content.get("stat", {"number": "92%", "label": "Accuracy Achieved"})
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(2.2), Inches(6.33), Inches(4.0))
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_WHITE
            card.line.color.rgb = COLOR_EMERALD
            card.line.width = Pt(2.0)
            
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(2.2), Inches(6.33), Inches(0.2))
            band.fill.solid()
            band.fill.fore_color.rgb = COLOR_EMERALD
            band.line.fill.background()
            
            tf = card.text_frame
            tf.word_wrap = True
            p_num = tf.paragraphs[0]
            p_num.text = f"\n{stat_data.get('number', '100%')}"
            p_num.font.bold = True
            p_num.font.size = Pt(64)
            p_num.font.color.rgb = COLOR_EMERALD
            p_num.alignment = PP_ALIGN.CENTER
            
            p_label = tf.add_paragraph()
            p_label.text = f"\n{stat_data.get('label', 'RESULT FINDING').upper()}"
            p_label.font.bold = True
            p_label.font.size = Pt(16)
            p_label.font.color.rgb = COLOR_ROYAL_BLUE
            p_label.alignment = PP_ALIGN.CENTER

    # LAYOUT 6: CONCLUSION & TAKEAWAY (Gold theme)
    elif layout_type == "conclusion":
        bullets = content.get("bullets", ["M.Sc. Research conclusion point."])
        
        accent_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(0.15), Inches(4.0))
        accent_card.fill.solid()
        accent_card.fill.fore_color.rgb = COLOR_GOLD
        accent_card.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(1.8), Inches(2.2), Inches(10.0), Inches(4.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p_c = tf.paragraphs[0]
        p_c.text = "🎯 THESIS SUMMARY"
        p_c.font.bold = True
        p_c.font.size = Pt(18)
        p_c.font.color.rgb = COLOR_GOLD
        
        for bullet in bullets[:4]:
            p = tf.add_paragraph()
            p.text = f"✓ {bullet}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT_DARK
            p.space_after = Pt(10)

    # FALLBACK / LAYOUT 7: HIGH-CONTRAST BULLETS
    else:
        bullets = content.get("bullets", ["No text details provided."])
        
        accent_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(0.15), Inches(4.0))
        accent_card.fill.solid()
        accent_card.fill.fore_color.rgb = COLOR_TEAL
        accent_card.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(1.8), Inches(2.2), Inches(10.0), Inches(4.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        for idx, bullet in enumerate(bullets[:5]):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = f"• {bullet}"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT_DARK
            p.space_after = Pt(10)

# ==========================================
# 7. SIMPLIFIED DYNAMIC TEMPLATE MATCHMAKING
# ==========================================
def get_5_templates_from_main():
    folder_path = "master_templates/"
    templates_to_use = []
    
    if os.path.exists(folder_path):
        all_files = [f for f in os.listdir(folder_path) if f.endswith('.pptx')]
        if len(all_files) > 0:
            selected = random.sample(all_files, min(len(all_files), 5))
            for t in selected:
                with open(os.path.join(folder_path, t), "rb") as f:
                    templates_to_use.append((t, f.read()))
                    
    if len(templates_to_use) == 0:
        for i in range(5):
            fallback_prs = Presentation()
            fallback_prs.slide_width = Inches(13.333)
            fallback_prs.slide_height = Inches(7.5)
            
            fallback_prs.slide_layouts[1].background.fill.solid()
            fallback_prs.slide_layouts[1].background.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            
            fallback_stream = io.BytesIO()
            fallback_prs.save(fallback_stream)
            fallback_stream.seek(0)
            templates_to_use.append((f"Fallback_Design_{i+1}.pptx", fallback_stream.read()))
            
    return templates_to_use

def generate_presentations_advanced(ai_data, master_templates):
    output_files = {}
    
    for i, (temp_name, temp_bytes) in enumerate(master_templates):
        prs = Presentation(io.BytesIO(temp_bytes))
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        slides_data = ai_data['slides']
        
        for slide_idx, slide_content in enumerate(slides_data):
            slide_layout = prs.slide_layouts[1]
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Draw premium M.Sc. slide visual card styling
            draw_visual_thesis_slide(new_slide, slide_content)
            
            # Programmatic Transition setup
            layout_type = slide_content.get("layout_type", "bullets")
            if layout_type == "methodology":
                set_slide_transition_and_timing(new_slide, transition_type="push", auto_advance_sec=60)
            elif layout_type == "results":
                set_slide_transition_and_timing(new_slide, transition_type="wipe", auto_advance_sec=60)
            else:
                set_slide_transition_and_timing(new_slide, transition_type="fade", auto_advance_sec=60)
                
            notes_slide = new_slide.notes_slide
            notes_slide.notes_text_frame.text = slide_content.get('script', "No speech generated.")
            
        output_stream = io.BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        output_files[f"Design_{temp_name.replace('.pptx', '')}_{i+1}.pptx"] = output_stream.getvalue()
        
    return output_files

# ==========================================
# 8. STREAMLIT WEB FRONTEND (UI)
# ==========================================
st.title("🎓 M.Sc. Premium Thesis Defense PPT Generator")
st.write("Convert your thesis chapters or rough document files into highly structured, dynamic presentations. Focuses on **Clear Objectives, Multi-Color process maps, Results Tables, and strict 6x6 Text Conciseness** automatically!")

uploaded_file = st.file_uploader("Upload your rough Thesis Document (.docx or .pptx)", type=["docx", "pptx"])
template_uploads = st.file_uploader("Upload custom Master PPTX Templates (Optional)", type=["pptx"], accept_multiple_files=True)

if st.button("Generate Premium Presentations", type="primary"):
    if not api_key:
        st.error("Please add your Gemini API Key first.")
    elif not uploaded_file:
        st.error("Please upload your rough document first.")
    else:
        with st.spinner("🤖 AI is reading, structuring, and applying academic standard visual codes..."):
            try:
                # 1. Read files
                if uploaded_file.name.endswith('.docx'):
                    raw_text = read_word_file(uploaded_file)
                    ai_raw_data = ai_process_word_advanced(raw_text)
                else:
                    raw_slides = read_ppt_file(uploaded_file)
                    ai_raw_data = ai_process_ppt_advanced(raw_slides)
                
                # Robust Normalization of AI JSON
                ai_data = normalize_ai_data(ai_raw_data)
                recommended_source = ai_data.get("recommended_source", "slidesgo")
                
                # 2. Template management
                templates_to_use = []
                if template_uploads:
                    for temp in template_uploads:
                        templates_to_use.append((temp.name, temp.read()))
                else:
                    templates_to_use = get_5_templates_from_main()

                # 3. Generate dynamic visual slides
                generated_ppts = generate_presentations_advanced(ai_data, templates_to_use)
                
                # 4. Create ZIP for easy download
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
                    for filename, file_data in generated_ppts.items():
                        zip_file.writestr(filename, file_data)
                
                zip_buffer.seek(0)
                
                st.success(f"🎉 M.Sc. Premium presentations successfully generated! AI recommends **{recommended_source.upper()}** style [2]!")
                
                # LIVE INTERACTIVE MERMAID FLOWCHART PREVIEW (Render directly in Web UI!)
                if "mermaid_diagram" in ai_data:
                    st.subheader("📊 Generated Methodology Flowchart (Live Preview)")
                    st.markdown(f"```mermaid\n{ai_data['mermaid_diagram']}\n```")
                
                st.download_button(
                    label="📥 Download Generated Presentations (ZIP)",
                    data=zip_buffer,
                    file_name="academic_thesis_presentations.zip",
                    mime="application/zip"
                )
                
            except Exception as e:
                st.error(f"Error occurred: {str(e)}")

st.write("---")
st.subheader("💡 Recommended Free Template Resources [2]")
st.write("Hamro system le auto-match garne official sites haru yahan chan. Tapai le aafule manparayeko template download garera 'Upload custom Master templates' section ma upload garna pani saknu huncha [2]:")

col1, col2 = st.columns(2)

with col1:
    st.markdown("""
    *   [**Slidesgo**](https://slidesgo.com/) [2]
        *   *Niche:* Modern, highly visual, creative academic decks. Perfect for classroom lectures or assignments [2].
    *   [**PresentationGO**](https://www.presentationgo.com/) [2]
        *   *Niche:* Clean grids, infographics, data tables, and professional sequential timelines [2].
    *   [**SlidesCarnival**](https://www.slidescarnival.com/) [2]
        *   *Niche:* Clean, traditional, minimalist academic & formal thesis layout designs [2].
    """)

with col2:
    st.markdown("""
    *   [**SlideEgg**](https://www.slideegg.com/) [2]
        *   *Niche:* Excellent structured diagrammatic templates, perfect for technical research [2].
    *   [**Microsoft Create**](https://create.microsoft.com/) [2]
        *   *Niche:* Official, highly stable corporate/academic standard Microsoft PowerPoint templates [2].
    *   [**Canva Presentations**](https://www.canva.com/) [2]
        *   *Niche:* Modern, visually rich slides (Can export easily to `.pptx` for our engine) [2].
    """)
